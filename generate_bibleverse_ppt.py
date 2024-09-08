
"""
get_bible_verse_from_website.py

Generate a Bible verse PowerPoint presentation from a text file.

Copyright (c) 2023 Astra
Maintainer: Astra <astralee95@gmail.com>

Licensed under the MIT License. See the LICENSE file for more details.
"""

import subprocess
import sys
import os
import ctypes
import time
import threading

# Set the console encoding to UTF-8
if os.name == 'nt':
    import msvcrt
    msvcrt.setmode(sys.stdout.fileno(), os.O_BINARY)
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.support.ui import Select
from selenium.webdriver.common.by import By
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC

from webdriver_manager.chrome import ChromeDriverManager
from webdriver_manager.firefox import GeckoDriverManager
from webdriver_manager.microsoft import EdgeChromiumDriverManager

from bs4 import BeautifulSoup

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsdecls
from lxml import etree
from pptx.oxml.xmlchemy import OxmlElement

import bible_menu
import color_helper

# Function to convert numeric values to Chinese traditional words
def num_to_chinese(num):
    chinese_nums = ["零", "一", "二", "三", "四", "五", "六", "七", "八", "九"]
    chinese_tens = ["", "十", "二十", "三十", "四十", "五十", "六十", "七十", "八十", "九十"]
    chinese_hundreds = ["", "一百", "二百", "三百", "四百", "五百", "六百", "七百", "八百", "九百"]

    if num < 10:
        return chinese_nums[num]
    elif num < 20:
        return "十" + chinese_nums[num % 10]
    elif num < 100:
        return chinese_tens[num // 10] + chinese_nums[num % 10]
    else:
        hundred_digit = num // 100
        ten_digit = (num % 100) // 10
        unit_digit = num % 10

        result = chinese_hundreds[hundred_digit]

        if ten_digit == 0:
            if unit_digit != 0:
                result += chinese_nums[unit_digit]
        else:
            result += chinese_tens[ten_digit]
            if unit_digit != 0:
                result += chinese_nums[unit_digit]

        return result

def fetch_bible_verses(query_file, browser):
    result = []

    text_thread.start()

    # Read book name, chapter number, and verse range from the provided query file
    with open(query_file, "r", encoding="utf-8") as file:
        for line in file:
            book_chap_verse = line.strip().split(':') 
            book_name = ''.join(filter(str.isalpha, book_chap_verse[0]))
            if book_name == '':
                continue

            # Check if there is a valid chapter number
            chap_str = ''.join(filter(str.isdigit, book_chap_verse[0]))
            chap_num = int(chap_str) if chap_str else None
            if chap_num < 1:
                print(f"{book_name}:{chap_num} chap number is invalid (>= 1), skip generating\n")
                continue

            # Determine if the book is Old Testament or New Testament
            form_index = None
            chapter_limit = None

            # Check in Old Testament books
            for initial, full, chapters, verses in bible_menu.ot:
                if book_name in (initial, full):
                    book_name = initial
                    form_index = 1  # Use form[1] for Old Testament
                    chapter_limit = chapters  # Get the number of chapters
                    break

            # Check in New Testament books if not found in Old Testament
            if form_index is None:
                for initial, full, chapters, verses in bible_menu.nt:
                    if book_name in (initial, full):
                        book_name = initial
                        form_index = 2  # Use form[2] for New Testament
                        chapter_limit = chapters  # Get the number of chapters

            if len(book_chap_verse) == 1: # use all verse
                if form_index == 1: # find in ot
                    for initial, full, chapters, verses in bible_menu.ot:
                        if book_name in (initial, full):
                            book_name = initial
                            form_index = 1  # Use form[1] for Old Testament
                            chapter_limit = chapters  # Get the number of chapters
                            verse_range = [1, verses[chap_num]]
                            break
                else: # find in nt
                    for initial, full, chapters, verses in bible_menu.nt:
                        if book_name in (initial, full):
                            book_name = initial
                            form_index = 1  # Use form[1] for Old Testament
                            chapter_limit = chapters  # Get the number of chapters
                            verse_range = [1, verses[chap_num ]]
                            break
            else:
                verse_range = book_chap_verse[1].strip()
                if '-' in verse_range:
                    verse_range = [int(x) for x in verse_range.split('-')]
                else:
                    verse_range = [int(verse_range), int(verse_range)]
            
            if min(verse_range) < 1 or verse_range[0] > verse_range[1]:
                print(f"v{book_chap_verse[1].strip()} verse range is invalid, skip generating\n")
                continue

            # Handle the case where the book is not found
            if form_index is None:
                print(f"Book '{book_name}' is not found in either the Old Testament or New Testament.")
            else:
                # Check if the chapter number is within the limits
                if 1 <= chap_num <= chapter_limit:
                    if form_index == 1:
                        if verse_range[0] <= len(bible_menu.ot[form_index - 1][3]) and verse_range[1] <= len(bible_menu.ot[form_index - 1][3]):
                            pass
                    elif form_index == 2:
                        if verse_range[0] <= len(bible_menu.nt[form_index - 1][3]) and verse_range[1] <= len(bible_menu.nt[form_index - 1][3]):
                            pass
                        
                    # Locate the specific form based on the determined index
                    # form = browser.find_element(By.XPATH, f"//form[{form_index}]")  # Adjust the index as needed
                    # Wait until the form is present before interacting with it
                    form = WebDriverWait(browser, 10).until(
                        EC.presence_of_element_located((By.XPATH, f"//form[{form_index}]"))
                    )

                    # Find the select element by name within the specific form
                    select_book = Select(form.find_element(By.NAME, "na"))
                    select_book.select_by_value(book_name)

                    # Find the select element by name
                    select_chap = Select(form.find_element(By.NAME, "chap"))
                    if book_name == "詩":
                        select_chap.select_by_value(f"{chap_num}")  # Format chap_num as a three-digit string
                    else:
                        select_chap.select_by_value(f"{chap_num:03}")  # Format chap_num as a three-digit string

                    # Find the submit button by name and click it
                    submit_button = form.find_element(By.NAME, "submit1")
                    submit_button.click()  # Click the submit button

                    # Get the page source
                    page_source = browser.page_source

                    # Use BeautifulSoup to parse the HTML
                    soup = BeautifulSoup(page_source, "html.parser")

                    # Find all <ol> tags
                    ol_elements = soup.find_all('ol')

                    if form_index == 1:
                        full_book_name = next((full for initial, full, _, _ in bible_menu.ot if initial == book_name), "Unknown")
                    elif form_index == 2:
                        full_book_name = next((full for initial, full, _, _ in bible_menu.nt if initial == book_name), "Unknown")
                    
                    if full_book_name == "詩篇":
                        chap_num_chinese = num_to_chinese(chap_num) + '篇'
                    else:
                        chap_num_chinese = num_to_chinese(chap_num) + '章'
                        
                    verse = ''
                    if verse_range[0] == verse_range[1]:
                        verse = str(verse_range[0])
                    else:
                        verse = f"{verse_range[0]}-{verse_range[1]}"

                    # print(f"{full_book_name}{chap_num_chinese}{verse}節")
                    
                    title = f"{full_book_name}{chap_num_chinese}{verse}節"

                    result.append(title)

                    # Extract text from each <ol> and join them, filtering out empty lines
                    ol_texts = []
                    for i, ol in enumerate(ol_elements):
                        # Find all <li> items within the current <ol>
                        li_items = ol.find_all('li')
                        
                        # Number each <li> item
                        for j, li in enumerate(li_items):
                            text = li.get_text(strip=True)  # Get the text of the <li> and strip whitespace
                            if text and verse_range[0] <= j + 1 <= verse_range[1]:  # Only add text within the verse range
                                ol_texts.append(f"{j + 1}.{text.replace('神', ' 神')}")  # Format with index

                        # Join all the texts from <ol> elements into a single string
                        final_text = "\n".join(ol_texts) + '\n'  # Single newline between different items

                        result.append(final_text)
                        # Print the text from <ol> tags
                        # print(final_text)

                        # Close the browser
                        browser.back()
                        print(f"> completed request: {book_name}{chap_num}:{verse_range[0]}-{verse_range[1]}\n")
                else:
                    print(f"Chapter number {chap_num} is outside the chapter limits for the book '{book_name}'.")

    # Save the text to output.txt file
    if result:
        # Convert the list elements to strings before writing to the file
        result_strings = [str(item) for item in result]

        # Write the list elements to the output.txt file
        with open("output.txt", "w", encoding="utf-8") as file:
            file.write("\n".join(result_strings))

    # Signal that fetching is complete
    stop_event.set()
    
    # Wait for the blinking text simulation to complete
    text_thread.join()

    return result  # Return the result list

######

# Define ANSI escape codes for blink and reset
BLINK = '\033[5m'
RESET = '\033[0m'

# Create a global event to signal stopping the blinking
stop_event = threading.Event()

# Simulate blinking text
def simulate_blinking_text(text, duration=1, interval=0.5):
    end_time = time.time() + duration
    while time.time() < end_time and not stop_event.is_set():
        print(f"{BLINK}{text}{RESET}", end='\r', flush=True)
        time.sleep(interval)
        print(f"{' ' * len(text)}", end='\r', flush=True)
        time.sleep(interval)
    print(f"{text}{RESET}")  # Print the text one last time without blinking

def add_text_shadow(run):
    shadow_xml = """
    <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="38100" dist="38100" dir="5400000" algn="ctr" rotWithShape="0">
            <a:srgbClr val="000000">
                <a:alpha val="50000"/>
            </a:srgbClr>
        </a:outerShdw>
    </a:effectLst>
    """
    shadow_element = parse_xml(shadow_xml)
    run_element = run._r
    run_properties = run_element.get_or_add_rPr()
    run_properties.append(shadow_element)

def set_background_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color_helper.COLOR_MAP[color])

def generate_ppt_from_txt(txt_file, ppt_file="verse.pptx", fontsize=30, fontcolor='white', bgcolor='default'):
    custom_template = "template.pptx"
    prs = Presentation(custom_template)

    # Read the text file with UTF-8 encoding
    with open(txt_file, 'r', encoding='utf-8') as file:
        lines = file.read()

    if not lines:
        print("The text file is empty.")
        return

    for paragraph in lines.split("\n\n"):
        paragraph = paragraph.strip()
        for line in paragraph.split("\n")[1:]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use the title and content layout
            title_shape = slide.shapes.title
            title_shape.text = paragraph.split("\n")[0]  # Set the title for the slide
            title_shape.text_frame.paragraphs[0].font.size = Pt(fontsize)
            title_shape.text_frame.paragraphs[0].font.underline = True
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color_helper.COLOR_MAP[fontcolor])
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            for run in title_shape.text_frame.paragraphs[0].runs:
                add_text_shadow(run)  # Add shadow to each run in the paragraph

            # Check for the first number in the sentence and use it for numbering
            first_number = line.split('.')[0]
            if first_number:
                pass

            # This line is part of the content
            text_frame = slide.placeholders[1].text_frame
            p = text_frame.add_paragraph()

            if text_frame.text and text_frame.text[0] == '\n':
                text_frame.text = text_frame.text[1:].strip()
                for i in range(len(text_frame.paragraphs)):
                    text_frame.paragraphs[i].text = line.split('.')[1]  # Set the content line

                    text_frame.paragraphs[i].alignment = PP_ALIGN.JUSTIFY
                    text_frame.paragraphs[i].font.size = Pt(fontsize)
                    text_frame.paragraphs[i].level = 0
                    text_frame.paragraphs[i].font.bold = True
                    text_frame.paragraphs[i].font.color.rgb = RGBColor(*color_helper.COLOR_MAP[fontcolor])                    
                    for run in text_frame.paragraphs[i].runs:
                        add_text_shadow(run)  # Add shadow to each run in the paragraph
                    
                    # Manually adjust left margin and first line indent (if needed)
                    text_frame.paragraphs[i].margin_left = Inches(0.75)
                    text_frame.paragraphs[i].margin_right = Inches(0.75)
                    text_frame.paragraphs[i].space_before = Pt(12)  # Space before the bullet point
                    text_frame.paragraphs[i].space_after = Pt(12)   # Space after the bullet point
                    text_frame.paragraphs[i].indent = Pt(12)   # Space after the bullet point
                    
                    # Set numbered bullet format by modifying the XML directly
                    pPr = p._element.get_or_add_pPr()
                    buAutoNum = OxmlElement('a:buAutoNum')
                    buAutoNum.set('type', 'arabicPeriod')
                    buAutoNum.set('startAt', first_number)
                    text_frame.paragraphs[i]._pPr.insert(0, buAutoNum)

            set_background_color(slide, bgcolor)

    # Use title as ppt_file name if ppt_file is empty
    if not ppt_file.lower().endswith('.pptx'):
        ppt_file += '.pptx'

    # Get the absolute path of the ppt_file
    ppt_file_path = os.path.abspath(ppt_file)
    
    try:
        prs.save(ppt_file)
        print(f"{color_helper.GREEN}{color_helper.BOLD} >> Successfully saved PowerPoint presentation to {color_helper.BLUE}{color_helper.BOLD}{ppt_file_path}{color_helper.RESET}")
    except Exception as e:
        print(f"Error: Failed to save PowerPoint presentation. {e}")
        print(f"{color_helper.RED}{color_helper.BOLD}Error: Failed to save PowerPoint presentation. {e}{color_helper.RESET}")

######

def initialize_browser(browser_type):
    options = None
    if browser_type == "chrome":
        options = webdriver.ChromeOptions()
        options.add_argument("--disable-notifications")
        options.add_argument("blink-settings=imagesEnabled=false")
        options.add_argument("--headless")
        options.add_argument("--ignore-certificate-errors")
        options.add_argument("--allow-running-insecure-content")
        options.add_argument("--log-level=1")  # 3 = FATAL, 2 = ERROR, 1 = WARNING, 0 = INFO
        return webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=options)
    elif browser_type == "firefox":
        options = webdriver.FirefoxOptions()
        options.add_argument("--headless")  # Run in headless mode
        options.set_preference("devtools.console.stdout.content", False)
        return webdriver.Firefox(service=Service(GeckoDriverManager().install()), options=options)
    elif browser_type == "edge":
        options = webdriver.EdgeOptions()
        options.add_argument("--headless")  # Run in headless mode
        options.add_experimental_option("excludeSwitches", ["enable-logging"])
        return webdriver.Edge(service=Service(EdgeChromiumDriverManager().install()), options=options)

def check_chrome_installed():
    chrome_paths = [
        # "/path/to/chrome",  # Add the full path to Chrome executable here
        "/usr/bin/google-chrome",  # on Linux
        "C:/Program Files/Google/Chrome/Application/chrome.exe",  # on Windows
    ]
    for path in chrome_paths:
        if os.path.exists(path):
            return True
    return False

def check_firefox_installed():
    try:
        # Attempt to run Chrome
        subprocess.run(["firefox"], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return True
    except FileNotFoundError:
        return False
    
def check_edge_installed():
    try:
        # Attempt to run Edge
        subprocess.run(["msedge"], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return True
    except FileNotFoundError:
        return False

# Enable virtual terminal processing on Windows
def enable_virtual_terminal_processing():
    kernel32 = ctypes.windll.kernel32
    handle = kernel32.GetStdHandle(-11)  # STD_OUTPUT_HANDLE
    mode = ctypes.c_ulong()
    kernel32.GetConsoleMode(handle, ctypes.byref(mode))
    mode.value |= 0x0004  # ENABLE_VIRTUAL_TERMINAL_PROCESSING
    kernel32.SetConsoleMode(handle, mode)

if __name__ == "__main__":

        # Call the function to enable virtual terminal processing
    if os.name == 'nt':
        enable_virtual_terminal_processing()

    # Welcome page with ASCII art
    ascii_art = """
    __    _ __    __                             ___               __ 
   / /_  (_) /_  / /__ _   _____  _____________ |__ \ ____  ____  / /_
  / __ \/ / __ \/ / _ \ | / / _ \/ ___/ ___/ _ \__/ // __ \/ __ \/ __/
 / /_/ / / /_/ / /  __/ |/ /  __/ /  (__  )  __/ __// /_/ / /_/ / /_  
/_.___/_/_.___/_/\___/|___/\___/_/  /____/\___/____/ .___/ .___/\__/  
                                                  /_/   /_/           
"""
    print(f"{color_helper.YELLOW}{color_helper.BOLD}{ascii_art}{color_helper.RESET}")
    print(f"{color_helper.YELLOW}{color_helper.REVERSE}Welcome to the PowerPoint Generator!{color_helper.RESET}\n")

    # print(f"{color_helper.YELLOW}{color_helper.BOLD}正在收集經文...{color_helper.RESET}\n")

    # Define the blinking text simulation
    text_thread = threading.Thread(target=simulate_blinking_text, args=(f"{color_helper.YELLOW}{color_helper.BOLD}正在收集經文...{color_helper.RESET}", 10, 0.5))

    # Check if Chrome, Firefox, or Edge is installed
    if check_chrome_installed():
        browser = initialize_browser("chrome")
    elif check_firefox_installed():
        browser = initialize_browser("firefox")
    elif check_edge_installed():
        browser = initialize_browser("edge")
    else:
        print("No supported browser is installed. Please install Chrome, Firefox, or Edge to use this script.")
        sys.exit(1)

    browser.get("https://springbible.fhl.net/Bible2/cgic201/read100.html")

    fetch_bible_verses(query_file="query_bible.txt", browser=browser)

    print(f"{color_helper.YELLOW}{color_helper.BOLD}\n正在產生 PowerPoint 簡報中 ...{color_helper.RESET}\n")
    
    generate_ppt_from_txt(txt_file="output.txt", ppt_file="verse.pptx", fontsize=55)
    
    print(f"{color_helper.YELLOW}{color_helper.BOLD}\nPowerPoint 簡報產生完畢。{color_helper.RESET}")
